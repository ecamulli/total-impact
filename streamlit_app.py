import streamlit as st
import requests
import pandas as pd
from datetime import datetime, timedelta
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed
import pytz
from pptx import Presentation
from pptx.util import Inches, Pt

@st.cache_data
def generate_excel_report(df, pivot, client_df, summary_client_df):
    output = BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        # Detailed Sensor Report
        df.to_excel(writer, sheet_name="Detailed Sensor Report", index=False)

        # Summary Sensor Report + Total
        pivot.to_excel(writer, sheet_name="Summary Sensor Report", index=False)
        ws1 = writer.sheets["Summary Sensor Report"]
        total_row_1 = len(pivot) + 1
        ws1.write(total_row_1, 0, "Total")
        ws1.write_formula(
            total_row_1, 4,
            f"=SUM(E2:E{total_row_1})",
            writer.book.add_format({"num_format": "0.00"})
        )

        # Detailed Client Report
        if not client_df.empty:
            client_df.to_excel(writer, sheet_name="Detailed Client Report", index=False)

        # Summary Client Report + Total
        if not summary_client_df.empty:
            summary_client_df.to_excel(writer, sheet_name="Summary Client Report", index=False)
            ws2 = writer.sheets["Summary Client Report"]
            total_row_2 = len(summary_client_df) + 1
            ws2.write(total_row_2, 0, "Total")
            avg_idx = summary_client_df.columns.get_loc("Avg Critical Hours Per Day")
            col_letter = chr(ord('A') + avg_idx)
            ws2.write_formula(
                total_row_2, avg_idx,
                f"=SUM({col_letter}2:{col_letter}{total_row_2})",
                writer.book.add_format({"num_format": "0.00"})
            )

        # Adjust column widths
        for sheet_name, data in {
            "Detailed Sensor Report": df,
            "Summary Sensor Report": pivot,
            "Detailed Client Report": client_df,
            "Summary Client Report": summary_client_df
        }.items():
            if not data.empty:
                worksheet = writer.sheets[sheet_name]
                for i, col in enumerate(data.columns):
                    worksheet.set_column(i, i, 23)

    output.seek(0)
    return output

@st.cache_data
def generate_ppt_summary(pivot, summary_client_df, account_name, from_str, to_str):
    prs = Presentation("New Impact Report Template.pptx")

    # Title slide (layout 0)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.placeholders[0].text = f"Impact Report for {account_name}"
    title_slide.placeholders[1].text = f"{from_str} to {to_str}"

    def add_table_slide(df, title):
        # Content slide (layout 2)
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = title
                break
        tbl = slide.shapes.add_table(
            rows=df.shape[0] + 1,
            cols=df.shape[1],
            left=Inches(0.5), top=Inches(1.5),
            width=Inches(9), height=Inches(0.3 * df.shape[0])
        ).table
        # header row
        for c, col_name in enumerate(df.columns):
            tbl.cell(0, c).text = str(col_name)
        # data rows
        for r, row in enumerate(df.values, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(10)

    add_table_slide(pivot.head(10), "📊 Summary Sensor Report")
    if not summary_client_df.empty:
        add_table_slide(summary_client_df.head(10), "👥 Summary Client Report")

    ppt_output = BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
    return ppt_output

# Main App
st.set_page_config(page_title="7SIGNAL Total Impact Report")
st.title("📊 7SIGNAL Total Impact Report")

# Inputs
account_name    = st.text_input("Account Name")
client_id       = st.text_input("Client ID")
client_secret   = st.text_input("Client Secret", type="password")
kpi_codes_input = st.text_input("Enter up to 4 sensor KPI codes (comma-separated)")

# Time range
st.markdown("### ⏱️ Select Date and Time Range (Eastern Time - ET)")
eastern       = pytz.timezone("US/Eastern")
now_et        = datetime.now(eastern)
default_start = now_et - timedelta(days=7)
from_date     = st.date_input("From Date",  value=default_start.date())
from_time     = st.time_input("From Time",  value=default_start.time())
to_date       = st.date_input("To Date",    value=now_et.date())
to_time       = st.time_input("To Time",    value=now_et.time())

from_datetime = eastern.localize(datetime.combine(from_date, from_time))
to_datetime   = eastern.localize(datetime.combine(to_date,   to_time))

if to_datetime > now_et:
    st.warning("'To' time cannot be in the future.")
    to_datetime = now_et
if from_datetime > to_datetime:
    st.error("'From' must be before 'To'")
    st.stop()

days_back = round((to_datetime - from_datetime).total_seconds() / 86400, 2)
if days_back > 30:
    st.error("Range cannot exceed 30 days.")
    st.stop()

# Show selected days
st.markdown(f"🗓 Selected Range: **{days_back} days**")
from_ts = int(from_datetime.timestamp() * 1000)
to_ts   = int(to_datetime.timestamp()   * 1000)

# API Helpers
def authenticate(cid, secret):
    try:
        r = requests.post(
            "https://api-v2.7signal.com/oauth2/token",
            data={"client_id": cid, "client_secret": secret, "grant_type": "client_credentials"},
            headers={"Content-Type": "application/x-www-form-urlencoded"}
        )
        return r.json().get("access_token") if r.status_code == 200 else None
    except:
        return None

def safe_get(url, headers):
    try:
        r = requests.get(url, headers=headers)
        return r if r.status_code == 200 else None
    except:
        return None

def get_service_areas(headers):
    r = safe_get("https://api-v2.7signal.com/topologies/sensors/serviceAreas", headers)
    return r.json().get("results", []) if r else []

def get_networks(headers):
    r = safe_get("https://api-v2.7signal.com/networks/sensors", headers)
    return r.json().get("results", []) if r else []

def get_kpi_data(headers, sa, net, code, from_ts, to_ts, days_back):
    url = (
        f"https://api-v2.7signal.com/kpis/sensors/service-areas/{sa['id']}"
        f"?kpiCodes={code}&from={from_ts}&to={to_ts}"
        f"&networkId={net['id']}&averaging=ALL"
    )
    r = safe_get(url, headers)
    if not r:
        return []
    results = []
    for result in r.json().get("results", []):
        for band in ["measurements24GHz", "measurements5GHz", "measurements6GHz"]:
            for m in result.get(band, []):
                samples    = m.get("samples") or 0
                sla        = m.get("slaValue") or 0
                total_mins = (to_ts - from_ts) / 1000 / 60
                crit_samp  = round(samples * (1 - sla/100), 2)
                crit_mins  = crit_samp * (total_mins / samples) if samples else 0
                results.append({
                    "Service Area": sa["name"],
                    "Network": net["name"],
                    "Band": {"measurements24GHz": "2.4GHz", "measurements5GHz": "5GHz", "measurements6GHz": "6GHz"}[band],
                    "Days Back": days_back,
                    "KPI Code": result.get("kpiCode"),
                    "KPI Name": result.get("name"),
                    "Samples": samples,
                    "SLA Value": sla,
                    "KPI Value": m.get("kpiValue"),
                    "Status": m.get("status"),
                    "Target Value": m.get("targetValue"),
                    "Critical Samples": crit_samp,
                    "Critical Hours Per Day": round(min(crit_mins/60/days_back, 24), 2)
                })
    return results

# Generate and download reports
if st.button("Generate Report!"):
    if not all([account_name, client_id, client_secret, kpi_codes_input]):
        st.warning("All fields are required.")
        st.stop()

    token = authenticate(client_id, client_secret)
    if not token:
        st.error("Authentication failed.")
        st.stop()
    
    headers = {"Authorization": f"Bearer {token}"}
    service_areas = get_service_areas(headers)
    networks      = get_networks(headers)
    kpi_codes     = [k.strip() for k in kpi_codes_input.split(",")][:4]

    results = []
    with ThreadPoolExecutor(max_workers=6) as ex:
        futures = [ex.submit(get_kpi_data, headers, sa, net, code, from_ts, to_ts, days_back)
                   for sa in service_areas for net in networks for code in kpi_codes]
        for f in as_completed(futures):
            results.extend(f.result())

    if not results:
        st.warning("No KPI data found.")
        st.stop()

    df = pd.DataFrame(results)
    pivot = (
        df.groupby(["Service Area", "Network", "Band"])['Critical Hours Per Day']
          .mean()
          .reset_index()
          .sort_values(by="Critical Hours Per Day", ascending=False)
    )
    pivot.insert(1, "Days Back", round(days_back,2))
    pivot["Critical Hours Per Day"] = pivot["Critical Hours Per Day"].round(2)
    pivot = pivot.rename(columns={"Critical Hours Per Day": "Avg Critical Hours Per Day"})

    client_url = (
        f"https://api-v2.7signal.com/kpis/agents/locations?from={from_ts}&to={to_ts}&includeClientCount=true"
    )
    r = safe_get(client_url, headers)
    rows = []
    if r:
        for loc in r.json().get('results', []):
            for t in loc.get('types', []):
                rows.append({
                    'Location': loc.get('locationName'),
                    'Client Count': loc.get('clientCount'),
                    'Days Back': round(days_back,2),
                    'Type': t.get('type').replace('_',' ').title(),
                    'Critical Hours Per Day': round(min((t.get('criticalSum') or 0)/60/days_back,24),2)
                })
    client_df = pd.DataFrame(rows)

    summary_client_df = pd.DataFrame()
    if not client_df.empty:
        summary_client_df = client_df.pivot_table(
            index=['Location','Client Count'], columns='Type',
            values='Critical Hours Per Day', aggfunc='mean'
        ).reset_index()
        summary_client_df.insert(1,'Days Back',round(days_back,2))
        type_cols = [c for c in summary_client_df.columns if c not in ['Location','Client Count','Days Back']]
        summary_client_df[type_cols] = summary_client_df[type_cols].round(2).fillna(0)
        summary_client_df['Avg Critical Hours Per Day'] = summary_client_df[type_cols].mean(axis=1).round(2)

    # Finally, download buttons
    excel_output = generate_excel_report(df, pivot, client_df, summary_client_df)
    from_str = from_datetime.strftime('%Y-%m-%d')
    to_str   = to_datetime.strftime('%Y-%m-%d')
    ppt_output = generate_ppt_summary(pivot, summary_client_df, account_name, from_str, to_str)

    st.download_button(
        "🗕 Download Excel Report",
        data=excel_output,
        file_name=f"{account_name}_impact_report_{from_str}_to_{to_str}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    st.download_button(
        "🎮 Download PowerPoint Summary",
        data=ppt_output,
        file_name=f"{account_name}_impact_report_{from_str}_to_{to_str}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
